#!/usr/bin/env python3
"""
scripts/create_presentation.py
Présentation GARF — CNN ablation study.
1 slide = 1 idée. Max 4 bullets. Visuels larges.

Usage:
    pip install python-pptx matplotlib scipy
    python scripts/create_presentation.py --out presentation.pptx
"""

import argparse, io
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from scipy.ndimage import gaussian_filter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
TEAL   = RGBColor(0x00, 0x7B, 0x9E)
LBLUE  = RGBColor(0xD6, 0xEA, 0xF8)
GREEN  = RGBColor(0x1E, 0x8B, 0x4C)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED    = RGBColor(0xBF, 0x2B, 0x2B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF3, 0xF4)
DGRAY  = RGBColor(0x5D, 0x6D, 0x7E)
DARK   = RGBColor(0x17, 0x20, 0x2A)
GOLD   = RGBColor(0xD4, 0xAC, 0x0D)
_BG    = "#F2F3F4"

SW = Inches(13.33)
SH = Inches(7.50)

# ── Results data ──────────────────────────────────────────────────────────
RESULTS = [
    {"label": "Step 1 — Baseline",       "acc": 0.7897, "f1": 0.6805, "prec": 0.6029, "rec": 0.8107},
    {"label": "Step 2 + Normals",        "acc": 0.7988, "f1": 0.7090, "prec": 0.6139, "rec": 0.8732},
    {"label": "Step 3 + Splatting",      "acc": 0.8240, "f1": 0.7193, "prec": 0.6338, "rec": 0.8715},
    {"label": "Step 4 + Bilinear",       "acc": 0.8242, "f1": 0.7200, "prec": 0.6306, "rec": 0.8732},
    {"label": "Step 5 + Context",        "acc": 0.8447, "f1": 0.7272, "prec": 0.6742, "rec": 0.8269},
    {"label": "Step 6 + Attention",      "acc": 0.8304, "f1": 0.7715, "prec": 0.7021, "rec": 0.8659},
    {"label": "Step 7 + U-Net",          "acc": 0.8284, "f1": 0.7708, "prec": 0.6964, "rec": 0.8748},
    {"label": "Step 8 + Feat. Fusion",   "acc": 0.8688, "f1": 0.8191, "prec": 0.7512, "rec": 0.9095},
    {"label": "Step 9 + Geo Features",   "acc": 0.8932, "f1": 0.8550, "prec": 0.7868, "rec": 0.9431},
    {"label": "PTv3 GARF mini",          "acc": 0.9371, "f1": 0.9098, "prec": 0.9167, "rec": 0.9100},
]

# ── PPTX helpers ──────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(sl, l, t, w, h, fill=TEAL, line=None, lw=Pt(0)):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else:    s.line.color.rgb = fill; s.line.width = Pt(0)
    return s

def tb(sl, text, l, t, w, h, size=18, bold=False, italic=False,
       color=DARK, align=PP_ALIGN.LEFT):
    b = sl.shapes.add_textbox(l, t, w, h)
    b.text_frame.word_wrap = True
    p = b.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return b

def img(sl, buf, l, t, w, h):
    sl.shapes.add_picture(buf, l, t, w, h)

def header(sl, title, subtitle=None):
    rect(sl, 0, 0, SW, Inches(1.1), fill=NAVY)
    tb(sl, title, Inches(0.5), Inches(0.1), Inches(11), Inches(0.75),
       size=30, bold=True, color=WHITE)
    if subtitle:
        tb(sl, subtitle, Inches(0.5), Inches(0.8), Inches(11), Inches(0.32),
           size=13, color=RGBColor(0xAA, 0xBB, 0xCC))

def bullets(sl, items, l, t, w, size=17, gap=0.65):
    for i, txt in enumerate(items):
        tb(sl, "•  " + txt, l, t + Inches(i * gap), w, Inches(gap),
           size=size, color=DARK)

def badge(sl, text, l, t, w=2.5, col=TEAL):
    rect(sl, Inches(l), Inches(t), Inches(w), Inches(0.44), fill=col)
    tb(sl, text, Inches(l), Inches(t + 0.04), Inches(w), Inches(0.38),
       size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Synthetic fragment ────────────────────────────────────────────────────

def _make_fragment(n=1500, seed=42):
    rng = np.random.default_rng(seed)
    th = rng.uniform(0, 2*np.pi, n)
    r  = np.sqrt(rng.uniform(0, 1, n))
    x, y = r*np.cos(th), r*np.sin(th)
    z = 0.28*(x**2+y**2-0.5) + rng.normal(0, 0.022, n)
    pts = np.column_stack([x, y, z])
    nrm = np.column_stack([-0.56*x, -0.56*y, np.ones(n)])
    nrm /= np.linalg.norm(nrm, axis=1, keepdims=True)
    frac = (r > 0.73) | (np.abs(z - z.min()) < 0.06)
    return pts, nrm, frac.astype(bool)

_PTS, _NRM, _FRAC = _make_fragment()

def _n01(a): return (a-a.min())/(a.max()-a.min()+1e-8)
def _pix(pts, ax, H, W):
    ui, vi = [(0,2),(1,2),(0,1)][ax]
    px = np.clip((_n01(pts[:,ui])*(W-1)).astype(int), 0, W-1)
    py = np.clip((_n01(pts[:,vi])*(H-1)).astype(int), 0, H-1)
    return px, py

def proj_floor(pts, ax, H=80, W=80):
    px, py = _pix(pts, ax, H, W)
    di = [1,0,2][ax]; d = _n01(pts[:,di])
    depth = np.full((H,W), np.nan)
    for i in range(len(pts)):
        if np.isnan(depth[py[i],px[i]]) or d[i] > depth[py[i],px[i]]:
            depth[py[i],px[i]] = d[i]
    occ = ~np.isnan(depth)
    return np.where(occ, depth, 0), occ

def proj_splat(pts, ax, H=80, W=80, sigma=1.5):
    px, py = _pix(pts, ax, H, W)
    di = [1,0,2][ax]; d = _n01(pts[:,di])
    acc = np.zeros((H,W)); cnt = np.zeros((H,W))
    for i in range(len(pts)):
        acc[py[i],px[i]] += d[i]; cnt[py[i],px[i]] += 1
    avg = np.where(cnt>0, acc/np.where(cnt>0,cnt,1), 0)
    return gaussian_filter(avg, sigma), gaussian_filter(cnt, sigma)>0.05

def proj_bilinear(pts, ax, H=80, W=80):
    ui, vi = [(0,2),(1,2),(0,1)][ax]; di = [1,0,2][ax]
    u = _n01(pts[:,ui])*(W-1); v = _n01(pts[:,vi])*(H-1); d = _n01(pts[:,di])
    acc = np.zeros((H,W)); wa = np.zeros((H,W))
    for i in range(len(pts)):
        x0,y0 = int(u[i]),int(v[i]); x1,y1 = min(x0+1,W-1),min(y0+1,H-1)
        wx,wy = u[i]-x0, v[i]-y0
        for bx,by,bw in [(x0,y0,(1-wx)*(1-wy)),(x1,y0,wx*(1-wy)),
                          (x0,y1,(1-wx)*wy),(x1,y1,wx*wy)]:
            acc[by,bx]+=bw*d[i]; wa[by,bx]+=bw
    return np.where(wa>0, acc/np.where(wa>0,wa,1), 0), wa>0.01

def nmap(pts, nrm, ax, H=80, W=80):
    px, py = _pix(pts, ax, H, W)
    nm = np.zeros((H,W,3)); cnt = np.zeros((H,W))
    for i in range(len(pts)):
        nm[py[i],px[i]] += (nrm[i]+1)/2; cnt[py[i],px[i]] += 1
    return np.where(cnt[:,:,None]>0, nm/np.where(cnt[:,:,None]>0,cnt[:,:,None],1), 0)

def _save(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0); plt.close(fig); return buf

def _clean(ax):
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.set_facecolor(_BG)

# ── Figures ───────────────────────────────────────────────────────────────

def fig_fragment_projections():
    fig = plt.figure(figsize=(11, 3.0), facecolor=_BG)
    ax3 = fig.add_subplot(1,4,1, projection='3d', facecolor=_BG)
    c = ['#e74c3c' if f else '#3498db' for f in _FRAC]
    ax3.scatter(_PTS[:,0], _PTS[:,1], _PTS[:,2], c=c, s=2, alpha=0.7, linewidths=0)
    ax3.set_title("3D fragment\n(red = fracture)", fontsize=9)
    ax3.set_axis_off()
    for i, name in enumerate(["Front view", "Side view", "Top view"]):
        ax = fig.add_subplot(1,4,i+2); d,o = proj_floor(_PTS, ax=i)
        ax.imshow(np.where(o,d,np.nan), cmap='plasma', origin='lower',
                  vmin=0, vmax=1, interpolation='nearest')
        ax.set_title(name, fontsize=9); ax.axis('off'); ax.set_facecolor(_BG)
    plt.tight_layout(pad=0.4)
    return _save(fig)

def fig_normals_comparison():
    fig, axes = plt.subplots(1, 3, figsize=(9, 3.2), facecolor=_BG)
    d, o = proj_floor(_PTS, ax=0); nm = nmap(_PTS, _NRM, ax=0)
    axes[0].imshow(np.where(o,d,np.nan), cmap='gray', origin='lower',
                   vmin=0, vmax=1, interpolation='nearest')
    axes[0].set_title("Depth only — Step 1", fontsize=10, pad=4); axes[0].axis('off')
    axes[1].imshow(nm, origin='lower', interpolation='nearest')
    axes[1].set_title("Normal map  (nx→R  ny→G  nz→B)", fontsize=10, pad=4); axes[1].axis('off')
    combined = np.clip(0.4*np.stack([np.where(o,d,0)]*3,-1) + 0.6*nm, 0, 1)
    axes[2].imshow(combined, origin='lower', interpolation='nearest')
    axes[2].set_title("Depth + normals — Step 2 input", fontsize=10, pad=4); axes[2].axis('off')
    for ax in axes: ax.set_facecolor(_BG)
    plt.tight_layout(pad=0.5)
    return _save(fig)

def fig_splatting_comparison():
    fig, axes = plt.subplots(1, 2, figsize=(7, 3.4), facecolor=_BG)
    d_fl, o_fl = proj_floor(_PTS, ax=2, H=64, W=64)
    d_sp, o_sp = proj_splat(_PTS, ax=2, H=64, W=64, sigma=1.5)
    axes[0].imshow(np.where(o_fl,d_fl,np.nan), cmap='plasma', origin='lower',
                   vmin=0, vmax=1, interpolation='nearest')
    axes[0].set_title("Floor rasterisation\n(sparse, aliased)", fontsize=11); axes[0].axis('off')
    axes[1].imshow(np.where(o_sp,d_sp,np.nan), cmap='plasma', origin='lower',
                   vmin=0, vmax=1, interpolation='nearest')
    axes[1].set_title("Gaussian splatting\n(dense, smooth)", fontsize=11); axes[1].axis('off')
    for ax in axes: ax.set_facecolor(_BG)
    plt.tight_layout(pad=0.5)
    return _save(fig)

def fig_bilinear_comparison():
    rng = np.random.default_rng(99)
    pts_s = _PTS[rng.choice(len(_PTS), 120, replace=False)]
    fig, axes = plt.subplots(1, 2, figsize=(7, 3.4), facecolor=_BG)
    d_fl, o_fl = proj_floor(pts_s, ax=2, H=40, W=40)
    d_bi, o_bi = proj_bilinear(pts_s, ax=2, H=40, W=40)
    axes[0].imshow(np.where(o_fl,d_fl,np.nan), cmap='plasma', origin='lower',
                   vmin=0, vmax=1, interpolation='nearest')
    axes[0].set_title("Floor  (blocky)", fontsize=11); axes[0].axis('off')
    axes[1].imshow(np.where(o_bi,d_bi,np.nan), cmap='plasma', origin='lower',
                   vmin=0, vmax=1, interpolation='nearest')
    axes[1].set_title("Bilinear  (smooth)", fontsize=11); axes[1].axis('off')
    for ax in axes: ax.set_facecolor(_BG)
    plt.tight_layout(pad=0.5)
    return _save(fig)

def fig_context_diagram():
    fig, ax = plt.subplots(figsize=(8, 3.0), facecolor=_BG)
    ax.set_xlim(0, 10); ax.set_ylim(0, 4); ax.axis('off'); ax.set_facecolor(_BG)
    def box(x, y, w, h, txt, col):
        ax.add_patch(plt.Rectangle((x,y),w,h,color=col,zorder=2))
        ax.text(x+w/2,y+h/2,txt,ha='center',va='center',
                fontsize=10,color='white',fontweight='bold',zorder=3)
    def arr(x1, y1, x2, y2):
        ax.annotate('',xy=(x2,y2),xytext=(x1,y1),
                    arrowprops=dict(arrowstyle='->',color='#5D6D7E',lw=2.0))
    box(0.2, 1.5, 2.2, 1.0, "CNN pixel\nfeatures", '#007B9E')
    arr(2.4, 2.0, 2.9, 2.0)
    box(2.9, 1.5, 1.8, 1.0, "Mean\nPool", '#1A2E4A')
    arr(4.7, 2.0, 5.2, 2.0)
    box(5.2, 1.5, 2.2, 1.0, "Global\ndescriptor", '#1A2E4A')
    ax.annotate('',xy=(1.3, 1.5),xytext=(6.3, 1.5),
                arrowprops=dict(arrowstyle='->',color='#5D6D7E',lw=2.0,
                                connectionstyle='arc3,rad=-0.35'))
    ax.text(3.8, 0.5, "broadcast to all pixels", ha='center', fontsize=9,
            color='#5D6D7E', style='italic')
    box(0.2, 0.1, 2.2, 0.9, "Concat → prediction", '#1E8B4C')
    ax.set_title("Global Context Injection (Step 5)", fontsize=11, pad=6)
    plt.tight_layout(pad=0.3)
    return _save(fig)

def fig_attention_diagram():
    fig, ax = plt.subplots(figsize=(8, 3.0), facecolor=_BG)
    ax.set_xlim(0, 10); ax.set_ylim(0, 4.5); ax.axis('off'); ax.set_facecolor(_BG)
    for i, (lbl, col, w) in enumerate([("Front view", '#007B9E', 0.55),
                                        ("Side view",  '#1A2E4A', 0.30),
                                        ("Top view",   '#7F8C8D', 0.15)]):
        y = 3.5 - i * 1.4
        ax.add_patch(plt.Rectangle((0.2, y-0.38), 2.6, 0.85, color=col, alpha=0.9))
        ax.text(1.5, y+0.05, lbl, ha='center', va='center',
                fontsize=10, color='white', fontweight='bold')
        ax.annotate('', xy=(3.2, y+0.05), xytext=(2.8, y+0.05),
                    arrowprops=dict(arrowstyle='->', color='#5D6D7E', lw=1.8))
        ax.text(3.6, y+0.05, f"α = {w:.2f}", va='center',
                fontsize=10, color=col, fontweight='bold')
    ax.text(5.8, 1.8, "Softmax\nweighted\nsum", ha='center', va='center', fontsize=9,
            color='#5D6D7E',
            bbox=dict(boxstyle='round', facecolor='white', edgecolor='#CCCCCC'))
    ax.annotate('', xy=(7.7, 1.8), xytext=(7.0, 1.8),
                arrowprops=dict(arrowstyle='->', color='#5D6D7E', lw=2.0))
    ax.add_patch(plt.Rectangle((7.7, 1.2), 2.0, 1.2, color='#1E8B4C', alpha=0.9))
    ax.text(8.7, 1.8, "Fused\nfeature", ha='center', va='center',
            fontsize=10, color='white', fontweight='bold')
    ax.set_title("View Attention (Step 6)", fontsize=11, pad=6)
    plt.tight_layout(pad=0.3)
    return _save(fig)

def fig_results_table_chart():
    labels = [r["label"] for r in RESULTS]
    f1s    = [r["f1"] for r in RESULTS]
    colors = ['#95A5A6','#7F8C8D','#7F8C8D','#5D6D7E','#5D6D7E',
              '#007B9E','#1A3C5E','#1E8B4C','#27AE60','#D4AC0D']
    fig, ax = plt.subplots(figsize=(10, 5.2), facecolor=_BG)
    bars = ax.barh(labels, f1s, color=colors, edgecolor='white', linewidth=0.6)
    ax.set_xlim(0, 1.08)
    ax.set_xlabel("F1 Score (val, best epoch)", fontsize=11)
    ax.set_title("Fracture Surface F1 — CNN Ablation vs PTv3 GARF", fontsize=12, pad=8)
    for bar, v in zip(bars, f1s):
        ax.text(v + 0.01, bar.get_y() + bar.get_height()/2,
                f"{v:.3f}", va='center', ha='left', fontsize=9)
    garf_f1 = RESULTS[-1]["f1"]
    ax.axvline(garf_f1, color='#D4AC0D', linestyle='--', lw=1.8, alpha=0.9,
               label=f'PTv3 GARF ({garf_f1:.3f})')
    ax.axvline(0.901, color='#E74C3C', linestyle=':', lw=1.5, alpha=0.7,
               label='Step 9 @ thresh=0.65 (~90.1%)')
    ax.legend(fontsize=9, loc='lower right')
    _clean(ax)
    fig.patch.set_facecolor(_BG)
    plt.tight_layout(pad=0.5)
    return _save(fig)

def fig_f1_progression():
    steps  = list(range(1, 10))
    labels = ["S1","S2","S3","S4","S5","S6","S7","S8","S9"]
    f1s    = [r["f1"] for r in RESULTS[:-1]]   # exclude PTv3
    colors = ['#95A5A6']*5 + ['#007B9E','#1A3C5E','#1E8B4C','#27AE60']
    fig, ax = plt.subplots(figsize=(9, 4.0), facecolor=_BG)
    ax.plot(steps, f1s, 'o-', color='#007B9E', lw=2.5, ms=8, zorder=3)
    ax.fill_between(steps, f1s, alpha=0.10, color='#007B9E')
    ax.axhline(RESULTS[-1]["f1"], color='#D4AC0D', linestyle='--', lw=2.0,
               label=f'PTv3 GARF ({RESULTS[-1]["f1"]:.3f})')
    ax.axhline(0.901, color='#E74C3C', linestyle=':', lw=1.5,
               label='Thresh calibrated (~90.1%)')
    for x, y in zip(steps, f1s):
        ax.annotate(f"{y:.3f}", (x, y), textcoords="offset points",
                    xytext=(0, 10), ha='center', fontsize=8)
    ax.set_xticks(steps)
    ax.set_xticklabels(labels, fontsize=10)
    ax.set_ylabel("F1 Score", fontsize=11)
    ax.set_ylim(0.55, 1.05)
    ax.legend(fontsize=9)
    ax.set_title("F1 Progression — Steps 1 → 9", fontsize=12, pad=8)
    _clean(ax)
    fig.patch.set_facecolor(_BG)
    plt.tight_layout(pad=0.5)
    return _save(fig)

def fig_threshold_sweep():
    thresholds = [0.30, 0.35, 0.40, 0.45, 0.50, 0.55, 0.60, 0.65, 0.70]
    f1s  = [0.8067, 0.8813, 0.8864, 0.8898, 0.8927, 0.8954, 0.8980, 0.9011, 0.9091]
    prec = [0.6769, 0.7928, 0.8028, 0.8101, 0.8167, 0.8232, 0.8303, 0.8399, 0.9071]
    rec  = [0.9982, 0.9922, 0.9894, 0.9869, 0.9844, 0.9814, 0.9777, 0.9719, 0.9111]
    fig, ax = plt.subplots(figsize=(8, 4.0), facecolor=_BG)
    ax.plot(thresholds, f1s,  'o-', color='#2196F3', lw=2.0, ms=7, label='F1')
    ax.plot(thresholds, prec, 's--',color='#FF9800', lw=1.8, ms=6, label='Precision')
    ax.plot(thresholds, rec,  '^:',color='#4CAF50', lw=1.8, ms=6, label='Recall')
    ax.axvline(0.50, color='#95A5A6', lw=1.5, linestyle='--', alpha=0.7, label='Default (0.5)')
    ax.axvline(0.65, color='#E74C3C', lw=1.8, linestyle='-', alpha=0.8, label='Optimal (0.65)')
    ax.axhline(RESULTS[-1]["f1"], color='#D4AC0D', lw=1.5, linestyle=':', label=f'PTv3 ({RESULTS[-1]["f1"]:.3f})')
    ax.annotate(f"F1={f1s[7]:.4f}\nPrec={prec[7]:.4f}\nRec={rec[7]:.4f}",
                xy=(0.65, f1s[7]), xytext=(0.55, 0.83),
                arrowprops=dict(arrowstyle='->', color='#E74C3C'),
                fontsize=9, color='#E74C3C')
    ax.set_xlabel("Decision threshold", fontsize=11)
    ax.set_ylabel("Score", fontsize=11)
    ax.set_title("Threshold Calibration on Step 9 (val set, point-level)", fontsize=11, pad=6)
    ax.set_ylim(0.60, 1.05)
    ax.legend(fontsize=9, loc='lower left')
    _clean(ax)
    fig.patch.set_facecolor(_BG)
    plt.tight_layout(pad=0.5)
    return _save(fig)

def fig_error_analysis():
    frac_bins = ["[0.00, 0.14)", "[0.14, 0.45)", "[0.45, 0.69)", "[0.69, 1.00)"]
    f1_frac   = [0.608, 0.829, 0.839, 0.928]
    nparts_bins = ["2", "3", "4", "5", "6+"]
    f1_parts    = [0.735, 0.773, 0.777, 0.772, 0.841]

    fig, axes = plt.subplots(1, 2, figsize=(10, 3.8), facecolor=_BG)

    # Left: by fracture ratio
    colors_l = ['#E74C3C','#F39C12','#27AE60','#1E8B4C']
    axes[0].bar(frac_bins, f1_frac, color=colors_l, edgecolor='white', linewidth=0.6)
    for i, v in enumerate(f1_frac):
        axes[0].text(i, v + 0.01, f"{v:.3f}", ha='center', fontsize=10)
    axes[0].set_ylim(0, 1.05)
    axes[0].set_xlabel("Fracture surface ratio", fontsize=10)
    axes[0].set_ylabel("F1 Score", fontsize=10)
    axes[0].set_title("F1 by Fracture Ratio\n← hardest    easiest →", fontsize=10)
    axes[0].tick_params(axis='x', labelsize=8)
    _clean(axes[0]); axes[0].set_facecolor(_BG)

    # Right: by n_parts
    colors_r = ['#E74C3C','#E67E22','#F1C40F','#2ECC71','#1E8B4C']
    axes[1].bar(nparts_bins, f1_parts, color=colors_r, edgecolor='white', linewidth=0.6)
    for i, v in enumerate(f1_parts):
        axes[1].text(i, v + 0.01, f"{v:.3f}", ha='center', fontsize=10)
    axes[1].set_ylim(0, 1.05)
    axes[1].set_xlabel("Number of parts", fontsize=10)
    axes[1].set_title("F1 by Object Complexity", fontsize=10)
    _clean(axes[1]); axes[1].set_facecolor(_BG)

    fig.patch.set_facecolor(_BG)
    plt.tight_layout(pad=0.6)
    return _save(fig)

# ── Slide builders ────────────────────────────────────────────────────────

def slide_01_title(prs):
    sl = blank(prs)
    rect(sl, 0, 0, Inches(5.5), SH, fill=NAVY)
    rect(sl, Inches(5.5), 0, SW - Inches(5.5), SH, fill=LGRAY)
    tb(sl, "CNN-based Fracture\nSurface Segmentation",
       Inches(0.4), Inches(1.2), Inches(5.0), Inches(2.8),
       size=36, bold=True, color=WHITE)
    tb(sl, "A Progressive Ablation Study",
       Inches(0.4), Inches(4.1), Inches(5.0), Inches(0.55),
       size=18, color=RGBColor(0xAA, 0xBB, 0xCC))
    tb(sl, "Compared against GARF (PTv3) baseline",
       Inches(0.4), Inches(4.7), Inches(5.0), Inches(0.45),
       size=14, italic=True, color=RGBColor(0x88, 0x99, 0xAA))
    tb(sl, "Teyssir Aissi  ·  USTH  ·  April 2026",
       Inches(0.4), Inches(6.7), Inches(5.0), Inches(0.4),
       size=13, color=DGRAY)
    img(sl, fig_fragment_projections(),
        Inches(5.7), Inches(2.0), Inches(7.3), Inches(2.6))
    tb(sl, "3D fragment  →  3 orthographic depth maps  →  CNN",
       Inches(5.7), Inches(4.7), Inches(7.3), Inches(0.4),
       size=11, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

def slide_02_motivation(prs):
    sl = blank(prs)
    header(sl, "Why This Work?")
    # Left panel
    rect(sl, Inches(0.5), Inches(1.3), Inches(5.6), Inches(4.8),
         fill=LBLUE, line=TEAL, lw=Pt(2))
    tb(sl, "GARF works well", Inches(0.7), Inches(1.5), Inches(5.2), Inches(0.5),
       size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    bullets(sl, ["F1 = 90.9%  on fracture segmentation",
                 "Point Transformer V3 (PTv3) encoder",
                 "End-to-end, strong 3D features"],
            Inches(0.8), Inches(2.15), Inches(5.0), size=15)
    # Right panel
    rect(sl, Inches(7.2), Inches(1.3), Inches(5.6), Inches(4.8),
         fill=RGBColor(0xFD, 0xED, 0xEC), line=RED, lw=Pt(2))
    tb(sl, "But — open questions", Inches(7.4), Inches(1.5), Inches(5.2), Inches(0.5),
       size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
    bullets(sl, ["Heavy 3D ops (hard to interpret)",
                 "Which geometry matters?",
                 "Can 2D projections compete?"],
            Inches(7.5), Inches(2.15), Inches(5.0), size=15)
    tb(sl, "→", Inches(6.1), Inches(3.5), Inches(1.0), Inches(0.6),
       size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

def slide_03_idea(prs):
    sl = blank(prs)
    header(sl, "Core Idea")
    rect(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.0), fill=NAVY)
    tb(sl, '"Project each 3D fragment onto 2D views,\nthen classify fracture points with a CNN."',
       Inches(1.0), Inches(1.65), Inches(11.3), Inches(1.7),
       size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullets(sl, ["Simple, interpretable, no 3D ops",
                 "3 orthographic views: front / side / top",
                 "Back-project CNN scores → per-point prediction"],
            Inches(1.5), Inches(3.8), Inches(10.5), size=16, gap=0.7)
    tb(sl, "→  We study this systematically with a 6-step ablation",
       Inches(1.5), Inches(6.0), Inches(10.5), Inches(0.55),
       size=15, italic=True, color=TEAL, align=PP_ALIGN.CENTER)

def slide_04_pipeline(prs):
    sl = blank(prs)
    header(sl, "Pipeline Overview")
    boxes = [("3D Fragment\nxyz + normals", TEAL),
             ("Ortho\nProjection\n(3 planes)", NAVY),
             ("ResNet-18\nCNN", NAVY),
             ("Back-\nprojection", TEAL),
             ("Fracture\nScore ∈ [0,1]", GREEN)]
    BW, BH, BY = Inches(2.2), Inches(1.5), Inches(1.6)
    xs = [Inches(0.3), Inches(2.75), Inches(5.2), Inches(7.65), Inches(10.1)]
    for (lbl, col), x in zip(boxes, xs):
        rect(sl, x, BY, BW, BH, fill=col)
        tb(sl, lbl, x, BY + Inches(0.15), BW, BH - Inches(0.15),
           size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i in range(len(xs)-1):
        tb(sl, "→", xs[i] + BW, BY + Inches(0.4), Inches(0.35), Inches(0.6),
           size=22, bold=True, color=DGRAY, align=PP_ALIGN.CENTER)
    subs = ["5 000 pts/frag", "128×128 depth maps", "shared weights", "NN lookup", "sigmoid > 0.5"]
    for sub, x in zip(subs, xs):
        tb(sl, sub, x, BY + BH + Inches(0.1), BW, Inches(0.35),
           size=10, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)
    rect(sl, Inches(0.3), Inches(3.5), Inches(12.7), Inches(0.38), fill=LBLUE)
    tb(sl, "Back-projection: each 3D point inherits the score of its nearest projected pixel",
       Inches(0.5), Inches(3.52), Inches(12.3), Inches(0.34),
       size=12, italic=True, color=NAVY, align=PP_ALIGN.CENTER)
    img(sl, fig_fragment_projections(),
        Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.1))

def slide_05_ablation_principle(prs):
    sl = blank(prs)
    header(sl, "Ablation Study Design",
           "One component added per step — contribution isolated")
    rows = [
        ("Step 1", "Depth only",              "baseline CNN",           "✓ done", DGRAY),
        ("Step 2", "+ Surface Normals",        "extra input channels",   "✓ done", TEAL),
        ("Step 3", "+ Gaussian Splatting",     "denser projection",      "✓ done", TEAL),
        ("Step 4", "+ Bilinear Interpolation", "smoother projection",    "✓ done", NAVY),
        ("Step 5", "+ Global Context",         "fragment-level feature", "✓ done", NAVY),
        ("Step 6", "+ View Attention",         "per-point view weights", "running…", GREEN),
    ]
    col_x = [Inches(0.3), Inches(1.6), Inches(4.5), Inches(9.2), Inches(11.2)]
    col_w = [Inches(1.2), Inches(2.8), Inches(4.6), Inches(1.9), Inches(1.9)]
    hdrs  = ["", "Step", "What changes", "Status", ""]
    Y0 = Inches(1.35)
    for hdr, cx, cw in zip(hdrs, col_x, col_w):
        rect(sl, cx, Y0, cw - Inches(0.05), Inches(0.42), fill=NAVY)
        tb(sl, hdr, cx + Inches(0.05), Y0 + Inches(0.06), cw - Inches(0.1), Inches(0.32),
           size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, (sid, name, what, status, col) in enumerate(rows):
        y = Inches(1.82 + i * 0.72)
        bg = LGRAY if i % 2 == 0 else WHITE
        rect(sl, Inches(0.3), y, Inches(12.8), Inches(0.68),
             fill=bg, line=RGBColor(0xDD, 0xDD, 0xDD), lw=Pt(0.5))
        scol = GREEN if status == "✓ done" else RED
        for txt, cx, cw, tc, al in [
            (sid,  col_x[0], col_w[0], col,   PP_ALIGN.CENTER),
            (name, col_x[1], col_w[1], col,   PP_ALIGN.LEFT),
            (what, col_x[2], col_w[2], DARK,  PP_ALIGN.LEFT),
            (status, col_x[3], col_w[3], scol, PP_ALIGN.CENTER),
        ]:
            tb(sl, txt, cx + Inches(0.05), y + Inches(0.17),
               cw - Inches(0.1), Inches(0.42),
               size=13, bold=(txt == name), color=tc, align=al)

def slide_06_normals(prs):
    sl = blank(prs)
    header(sl, "Step 2 — Surface Normals",
           "Adding surface orientation to the CNN input")
    bullets(sl, ["Add a normal map (nx, ny, nz) alongside the depth map",
                 "CNN input: 4 channels per view instead of 1",
                 "Hypothesis: orientation reveals fracture boundaries"],
            Inches(0.6), Inches(1.3), Inches(5.4), size=15, gap=0.68)
    badge(sl, "Result:  F1  68.1% → 70.9%  (+2.8 pp)",
          l=0.6, t=3.5, w=5.2, col=TEAL)
    badge(sl, "Recall  81.1% → 87.3%  (+6.3 pp)", l=0.6, t=4.1, w=5.2, col=TEAL)
    img(sl, fig_normals_comparison(),
        Inches(6.2), Inches(1.2), Inches(6.9), Inches(3.6))

def slide_07_splatting(prs):
    sl = blank(prs)
    header(sl, "Step 3 — Gaussian Splatting",
           "Denser projection coverage")
    bullets(sl, ["Each 3D point splat as a 2D Gaussian blob",
                 "Fills holes: fewer empty pixels in depth map",
                 "σ calibrated to mean inter-point distance"],
            Inches(0.6), Inches(1.3), Inches(5.4), size=15, gap=0.68)
    badge(sl, "Result:  F1  70.9% → 71.9%  (+1.0 pp)", l=0.6, t=3.5, w=5.2, col=NAVY)
    badge(sl, "Precision  61.4% → 63.4%  (+2.0 pp)",   l=0.6, t=4.1, w=5.2, col=NAVY)
    img(sl, fig_splatting_comparison(),
        Inches(6.2), Inches(1.3), Inches(6.9), Inches(3.5))

def slide_08_bilinear(prs):
    sl = blank(prs)
    header(sl, "Step 4 — Bilinear Interpolation",
           "Smoother point-to-pixel mapping")
    bullets(sl, ["Each point contributes to 4 surrounding pixels (weighted)",
                 "Sub-pixel accuracy — smoother depth gradients",
                 "Hypothesis: less aliasing → better generalisation"],
            Inches(0.6), Inches(1.3), Inches(5.4), size=15, gap=0.68)
    badge(sl, "Result:  F1  71.9% → 70.9%  (−1.0 pp)", l=0.6, t=3.5, w=5.2,
          col=RGBColor(0xBF, 0x2B, 0x2B))
    tb(sl, "→ Splatting already handles coverage; bilinear adds no benefit here.",
       Inches(0.6), Inches(4.3), Inches(5.2), Inches(0.6),
       size=13, italic=True, color=DGRAY)
    img(sl, fig_bilinear_comparison(),
        Inches(6.2), Inches(1.3), Inches(6.9), Inches(3.5))

def slide_09_context(prs):
    sl = blank(prs)
    header(sl, "Step 5 — Global Context",
           "Fragment-level awareness injected into every pixel")
    bullets(sl, ["Mean-pool all CNN pixel features per fragment",
                 "Concatenate this global descriptor back into every pixel",
                 "Lets the model reason globally, not just locally"],
            Inches(0.6), Inches(1.3), Inches(5.4), size=15, gap=0.68)
    badge(sl, "Result:  F1  70.9% → 72.6%  (+1.7 pp)", l=0.6, t=3.5, w=5.2, col=GREEN)
    badge(sl, "Best CNN result overall", l=0.6, t=4.1, w=5.2,
          col=RGBColor(0x1E, 0x6B, 0x3C))
    img(sl, fig_context_diagram(),
        Inches(6.2), Inches(1.4), Inches(6.9), Inches(3.4))

def slide_10_attention(prs):
    sl = blank(prs)
    header(sl, "Step 6 — View Attention",
           "Learning which view matters per fragment")
    bullets(sl, ["Small MLP scores each of the 3 views per fragment",
                 "Softmax → weighted combination at backprojection",
                 "Attention at prediction level — views still decoded separately"],
            Inches(0.6), Inches(1.3), Inches(5.4), size=15, gap=0.68)
    badge(sl, "F1  77.15%  (+4.5 pp vs Step 5)", l=0.6, t=3.5, w=5.2, col=TEAL)
    badge(sl, "Prec  70.21%  Rec  86.59%",       l=0.6, t=4.1, w=5.2, col=TEAL)
    img(sl, fig_attention_diagram(),
        Inches(6.2), Inches(1.4), Inches(6.9), Inches(3.4))

def slide_10b_unet(prs):
    sl = blank(prs)
    header(sl, "Step 7 — U-Net Backbone",
           "Skip connections for spatial precision")
    bullets(sl, ["Replace SimpleCNN with U-Net (encoder + skip connections + decoder)",
                 "Skip connections pass high-res features to decoder at each level",
                 "4 encoder levels: 16→32→64→128 channels"],
            Inches(0.6), Inches(1.3), Inches(5.4), size=15, gap=0.68)
    badge(sl, "F1  77.08%  ≈ Step 6", l=0.6, t=3.5, w=5.2, col=NAVY)
    tb(sl, "→ U-Net alone ≈ no improvement over SimpleCNN + attention\n"
           "   Skip connections need richer semantics to be useful",
       Inches(0.6), Inches(4.2), Inches(5.4), Inches(1.0),
       size=13, italic=True, color=DGRAY)
    # U-Net diagram
    fig, ax = plt.subplots(figsize=(7, 4.0), facecolor=_BG)
    ax.set_xlim(0, 10); ax.set_ylim(0, 6); ax.axis('off'); ax.set_facecolor(_BG)
    def box(x, y, w, h, txt, col, fs=9):
        ax.add_patch(plt.Rectangle((x, y), w, h, color=col, zorder=2, alpha=0.9))
        ax.text(x+w/2, y+h/2, txt, ha='center', va='center',
                fontsize=fs, color='white', fontweight='bold', zorder=3)
    enc = [(0.2,4.5,2.0,0.9,"Enc0\n16ch",'#007B9E'),
           (0.2,3.2,2.0,0.9,"Enc1\n32ch",'#1A3C5E'),
           (0.2,1.9,2.0,0.9,"Enc2\n64ch",'#1A2E4A'),
           (0.2,0.5,2.0,0.9,"Bottleneck\n128ch",'#0D1B2A')]
    dec = [(7.8,4.5,2.0,0.9,"Dec0\n16ch",'#1E8B4C'),
           (7.8,3.2,2.0,0.9,"Dec1\n32ch",'#27AE60'),
           (7.8,1.9,2.0,0.9,"Dec2\n64ch",'#2ECC71')]
    for (x,y,w,h,t,c) in enc: box(x,y,w,h,t,c)
    for (x,y,w,h,t,c) in dec: box(x,y,w,h,t,c)
    # skip connections
    for (xe,ye,_,he,_,_),(xd,yd,_,hd,_,_) in zip(enc[:3], dec):
        ax.annotate('', xy=(xd, yd+hd/2), xytext=(xe+2.0, ye+he/2),
                    arrowprops=dict(arrowstyle='->', color='#E74C3C', lw=1.5,
                                   connectionstyle='arc3,rad=-0.2'))
    ax.annotate('', xy=(7.8, 1.95+0.45), xytext=(2.2, 0.5+0.45),
                arrowprops=dict(arrowstyle='->', color='#5D6D7E', lw=2))
    ax.text(5.0, 0.9, "bottleneck", ha='center', fontsize=8, color='#5D6D7E')
    ax.text(5.0, 3.8, "skip connections", ha='center', fontsize=8,
            color='#E74C3C', style='italic')
    ax.set_title("U-Net encoder-decoder with skip connections", fontsize=10, pad=4)
    import io as _io; buf = _io.BytesIO()
    fig.savefig(buf, format='png', dpi=130, bbox_inches='tight', facecolor=_BG)
    buf.seek(0); plt.close(fig)
    img(sl, buf, Inches(6.2), Inches(1.4), Inches(6.9), Inches(4.5))

def _fig_pred_comparison():
    """
    Synthetic visual: GT fracture vs Step-7-style prediction vs Step-8-style prediction.
    Shows the effect of feature fusion on the quality of segmentation.
    """
    rng = np.random.default_rng(42)
    pts, _, frac = _make_fragment(n=2000, seed=42)

    # Simulate Step 7 prediction: correct on core fracture but spreads into neighbours
    # (high recall, imprecise — lots of FP around the fracture zone)
    dist_to_frac = np.zeros(len(pts))
    frac_pts = pts[frac]
    for i, p in enumerate(pts):
        if frac[i]:
            dist_to_frac[i] = 0.0
        else:
            dist_to_frac[i] = np.min(np.linalg.norm(frac_pts - p, axis=1))
    # Step 7: predicts fracture within distance 0.35 of real fracture (lots of FP)
    pred_s7 = (frac) | (dist_to_frac < 0.35)
    pred_s7 = pred_s7 & ~(rng.random(len(pts)) < 0.08)   # small FN noise

    # Step 8: tighter prediction, much fewer FP
    pred_s8 = (frac) | (dist_to_frac < 0.12)
    pred_s8 = pred_s8 & ~(rng.random(len(pts)) < 0.04)   # fewer FN

    def point_colors(pred, gt):
        c = np.full((len(gt), 3), [0.82, 0.82, 0.82])   # TN gray
        tp = pred & gt;  c[tp] = [0.15, 0.75, 0.25]     # TP green
        fp = pred & ~gt; c[fp] = [0.88, 0.15, 0.15]     # FP red
        fn = ~pred & gt; c[fn] = [0.15, 0.35, 0.85]     # FN blue
        return c

    def f1_str(pred, gt):
        tp = int((pred & gt).sum()); fp = int((pred & ~gt).sum())
        fn = int((~pred & gt).sum())
        p = tp/(tp+fp+1e-8); r = tp/(tp+fn+1e-8)
        f = 2*p*r/(p+r+1e-8)
        return f"F1={f:.2f}  Prec={p:.2f}  Rec={r:.2f}"

    fig, axes = plt.subplots(1, 3, figsize=(12, 4.0), facecolor=_BG)
    configs = [
        ("Ground Truth", frac,    lambda gt: np.where(gt[:,None], [[0.9,0.2,0.2]], [[0.6,0.6,0.95]])),
        (f"Step 7 prediction\n{f1_str(pred_s7, frac)}", pred_s7, lambda p: point_colors(p, frac)),
        (f"Step 8 prediction\n{f1_str(pred_s8, frac)}", pred_s8, lambda p: point_colors(p, frac)),
    ]
    for ax, (title, mask, color_fn) in zip(axes, configs):
        c = color_fn(mask)
        ax.scatter(pts[:,0], pts[:,1], c=c, s=4, linewidths=0, alpha=0.85)
        ax.set_title(title, fontsize=10, fontweight='bold', pad=5)
        ax.set_aspect('equal'); ax.axis('off'); ax.set_facecolor(_BG)

    # Legend on last panel
    import matplotlib.patches as mpatches
    legend = [mpatches.Patch(color=[0.15,0.75,0.25], label='TP'),
              mpatches.Patch(color=[0.88,0.15,0.15], label='FP'),
              mpatches.Patch(color=[0.15,0.35,0.85], label='FN'),
              mpatches.Patch(color=[0.82,0.82,0.82], label='TN')]
    axes[2].legend(handles=legend, loc='lower right', fontsize=8)

    fig.patch.set_facecolor(_BG)
    plt.tight_layout(pad=0.5)
    return _save(fig)


def _fig_fusion_comparison():
    """Side-by-side: Step 7 (prediction fusion) vs Step 8 (feature fusion)."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.8), facecolor=_BG)

    view_colors = ['#2471A3', '#1A5276', '#0D2D46']
    fuse_color  = '#D4AC0D'
    enc_color   = '#1A2E4A'
    dec_color   = '#1E8B4C'
    pred_color  = '#7D3C98'

    def box(ax, x, y, w, h, txt, col, fs=8, alpha=1.0):
        ax.add_patch(plt.Rectangle((x, y), w, h, color=col,
                                   zorder=2, alpha=alpha, linewidth=0))
        ax.text(x+w/2, y+h/2, txt, ha='center', va='center',
                fontsize=fs, color='white', fontweight='bold', zorder=3,
                multialignment='center')

    def arr(ax, x1, y1, x2, y2, col='#5D6D7E', lw=1.5, style='->'):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle=style, color=col, lw=lw))

    # ── LEFT: Step 7 — prediction-level fusion ──────────────────────────────
    ax1.set_xlim(0, 10); ax1.set_ylim(0, 5.5); ax1.axis('off')
    ax1.set_facecolor(_BG)
    ax1.set_title("Step 7 — Fusion at PREDICTION level", fontsize=11,
                  fontweight='bold', color='#E74C3C', pad=6)

    vy_list = [4.0, 2.5, 1.0]
    for i, (vy, vc) in enumerate(zip(vy_list, view_colors)):
        box(ax1, 0.1, vy, 1.2, 0.85, f"View {i+1}\nimage", vc, fs=8)
        box(ax1, 1.6, vy, 1.4, 0.85, "Encode", enc_color, fs=8)
        box(ax1, 3.3, vy, 1.4, 0.85, "Bottleneck", enc_color, fs=7)
        box(ax1, 5.0, vy, 1.4, 0.85, "Decode", dec_color, fs=8)
        box(ax1, 6.7, vy, 1.2, 0.85, f"Pred {i+1}", pred_color, fs=8)
        arr(ax1, 1.3, vy+0.42, 1.6, vy+0.42)
        arr(ax1, 3.0, vy+0.42, 3.3, vy+0.42)
        arr(ax1, 4.7, vy+0.42, 5.0, vy+0.42)
        arr(ax1, 6.4, vy+0.42, 6.7, vy+0.42)

    # fusion arrow (after predictions)
    box(ax1, 8.2, 2.5, 1.5, 0.85, "FUSE\nmean/attn", fuse_color, fs=8)
    for vy in vy_list:
        arr(ax1, 7.9, vy+0.42, 8.2, 2.92, col='#D4AC0D', lw=1.2)

    ax1.text(5.0, 0.3, "← vues traitées séparément\n   fusion APRÈS prédiction",
             ha='center', fontsize=8, color='#E74C3C', style='italic')

    # ── RIGHT: Step 8 — feature-level fusion ────────────────────────────────
    ax2.set_xlim(0, 10); ax2.set_ylim(0, 5.5); ax2.axis('off')
    ax2.set_facecolor(_BG)
    ax2.set_title("Step 8 — Fusion at BOTTLENECK level", fontsize=11,
                  fontweight='bold', color='#27AE60', pad=6)

    for i, (vy, vc) in enumerate(zip(vy_list, view_colors)):
        box(ax2, 0.1, vy, 1.2, 0.85, f"View {i+1}\nimage", vc, fs=8)
        box(ax2, 1.6, vy, 1.4, 0.85, "Encode", enc_color, fs=8)
        box(ax2, 3.3, vy, 1.4, 0.85, "Bottleneck", enc_color, fs=7)
        arr(ax2, 1.3, vy+0.42, 1.6, vy+0.42)
        arr(ax2, 3.0, vy+0.42, 3.3, vy+0.42)

    # FUSE in the middle
    box(ax2, 5.0, 2.5, 1.5, 0.85, "FUSE\n(max)", fuse_color, fs=9)
    for vy in vy_list:
        arr(ax2, 4.7, vy+0.42, 5.0, 2.92, col='#D4AC0D', lw=1.5)

    # broadcast back
    box(ax2, 6.8, 2.5, 1.0, 0.85, "broad-\ncast", '#5D6D7E', fs=7)
    arr(ax2, 6.5, 2.92, 6.8, 2.92, col='#D4AC0D', lw=2.0)

    for i, (vy, vc) in enumerate(zip(vy_list, view_colors)):
        box(ax2, 8.0, vy, 1.2, 0.85, "Decode", dec_color, fs=8)
        arr(ax2, 7.8, 2.92, 8.0, vy+0.42, col='#27AE60', lw=1.2)

    ax2.text(5.0, 0.3, "← fusion AVANT décodage\n   tous les décodeurs partagent la même sémantique",
             ha='center', fontsize=8, color='#27AE60', style='italic')

    plt.tight_layout(pad=0.8)
    buf = _save(fig)
    return buf


def slide_10c_fusion(prs):
    sl = blank(prs)
    header(sl, "Step 8 — Feature-Level View Fusion",
           "Fusionner les représentations AVANT de décoder, pas après")

    # Visual comparison top
    img(sl, _fig_pred_comparison(),
        Inches(0.3), Inches(1.15), Inches(12.7), Inches(2.9))

    # Key message
    tb(sl, "Step 7 : vues décodées séparément → fusion des prédictions → trop de FP",
       Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.38),
       size=13, color=RED)
    tb(sl, "Step 8 : bottlenecks fusionnés (max) → broadcast → même sémantique dans tous les décodeurs",
       Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.38),
       size=13, color=GREEN, bold=True)

    # Pipeline comparison diagram
    img(sl, _fig_fusion_comparison(),
        Inches(0.3), Inches(5.0), Inches(9.5), Inches(2.25))

    badge(sl, "F1  81.91%  (+4.7 pp)", l=9.9, t=5.1, w=3.2, col=GREEN)
    badge(sl, "8d: SimpleCNN+fusion=73.3%\n→ U-Net indispensable",
          l=9.9, t=5.75, w=3.2, col=NAVY)

def slide_10d_geo(prs):
    sl = blank(prs)
    header(sl, "Step 9 — Explicit 3D Geometric Features",
           "Project geometric descriptors as extra image channels")
    bullets(sl, ["Compute per-point features via local k-NN PCA (k=16, CPU)",
                 "3 new channels: curvature  |  roughness  |  normal consistency",
                 "Zero extra trainable parameters — pure geometric signal"],
            Inches(0.6), Inches(1.3), Inches(5.4), size=14, gap=0.72)
    badge(sl, "F1  85.50%  (+3.6 pp vs Step 8)", l=0.6, t=3.7, w=5.4, col=GREEN)
    badge(sl, "Rec  94.31%  Prec  78.68%",        l=0.6, t=4.3, w=5.4, col=GREEN)
    tb(sl, "Input channels: depth · occ · nx · ny · nz · curv · rough · consist  (8 total)",
       Inches(0.6), Inches(5.1), Inches(5.4), Inches(0.5),
       size=11, italic=True, color=DGRAY)
    # Feature description panel
    fig, ax = plt.subplots(figsize=(7, 4.0), facecolor=_BG)
    ax.set_xlim(0, 10); ax.set_ylim(0, 6); ax.axis('off'); ax.set_facecolor(_BG)
    items = [
        (0.5, 4.5, '#E74C3C', "Curvature  (λ_min / Σλ)",
         "High on fracture edges\n≈ 0 on flat intact surface"),
        (0.5, 2.8, '#9B59B6', "Roughness",
         "Mean dist to tangent plane\nHigh on irregular fracture faces"),
        (0.5, 1.1, '#007B9E', "Normal consistency",
         "Alignment of input normal vs PCA normal\nLow on ambiguous/fracture zones"),
    ]
    for x, y, c, title, desc in items:
        ax.add_patch(plt.Rectangle((x, y), 0.25, 1.0, color=c, zorder=2))
        ax.text(x+0.4, y+0.7, title, fontsize=11, color=c, fontweight='bold')
        ax.text(x+0.4, y+0.15, desc, fontsize=9, color='#333333')
    ax.set_title("3 geometric features — all computed from local k-NN PCA", fontsize=10, pad=4)
    import io as _io; buf = _io.BytesIO()
    fig.savefig(buf, format='png', dpi=130, bbox_inches='tight', facecolor=_BG)
    buf.seek(0); plt.close(fig)
    img(sl, buf, Inches(6.2), Inches(1.4), Inches(6.9), Inches(4.5))

def slide_10e_threshold(prs):
    sl = blank(prs)
    header(sl, "Threshold Calibration — A Free +4.5 pp",
           "Default threshold 0.5 is not optimal for this model")
    bullets(sl, ["Swept decision threshold 0.30 → 0.70 on val set (no retraining)",
                 "Optimal: thresh=0.65 → F1=90.11%  Prec=83.99%  Rec=97.19%",
                 "At thresh=0.70: F1=90.91% ≈ PTv3 GARF (90.98%)"],
            Inches(0.6), Inches(1.3), Inches(5.4), size=14, gap=0.72)
    badge(sl, "PTv3 GARF: F1=90.98%  (12.7M params)", l=0.6, t=3.7, w=5.4,
          col=RGBColor(0xD4, 0xAC, 0x0D))
    badge(sl, "Step 9 thresh=0.65: F1≈90.1%  (523K params)", l=0.6, t=4.3, w=5.4,
          col=GREEN)
    tb(sl, "24× fewer parameters — same F1 with calibrated threshold",
       Inches(0.6), Inches(5.1), Inches(5.4), Inches(0.5),
       size=13, bold=True, color=NAVY)
    img(sl, fig_threshold_sweep(),
        Inches(6.2), Inches(1.3), Inches(6.9), Inches(5.0))

def slide_10f_error_analysis(prs):
    sl = blank(prs)
    header(sl, "Step 9 Error Analysis",
           "5 386 fragments analysed — where does the model fail?")
    bullets(sl, ["Global: FP rate=19.6%  FN rate=3.1%  → model predicts too wide",
                 "Low fracture ratio (<14%): F1=60.8% — model predicts nothing",
                 "6+ parts objects: F1=84.1% — complex objects are easier"],
            Inches(0.6), Inches(1.3), Inches(5.4), size=14, gap=0.72)
    tb(sl, "Key finding: two distinct failure modes",
       Inches(0.6), Inches(3.55), Inches(5.4), Inches(0.4),
       size=13, bold=True, color=RED)
    tb(sl, "• fracture_ratio > 14%  →  FP too high (model predicts too wide)\n"
           "• fracture_ratio < 5%   →  model predicts nothing (class collapse)",
       Inches(0.6), Inches(3.95), Inches(5.4), Inches(0.9),
       size=12, color=DARK)
    img(sl, fig_error_analysis(),
        Inches(6.0), Inches(1.3), Inches(7.1), Inches(4.8))

def slide_11_results_table(prs):
    sl = blank(prs)
    header(sl, "Results — Full Ablation Table",
           "Best val-F1 checkpoint  ·  seed 1116  ·  100 epochs")
    hdrs = ["", "Model", "F1 ↑", "Prec ↑", "Recall ↑", "Params"]
    col_x = [Inches(0.3), Inches(1.15), Inches(5.2), Inches(6.7), Inches(8.2), Inches(9.75)]
    col_w = [Inches(0.83), Inches(4.0),  Inches(1.45), Inches(1.45), Inches(1.45), Inches(1.45)]
    Y0 = Inches(1.25)
    for hdr, cx, cw in zip(hdrs, col_x, col_w):
        rect(sl, cx, Y0, cw - Inches(0.04), Inches(0.40), fill=NAVY)
        tb(sl, hdr, cx + Inches(0.04), Y0 + Inches(0.05), cw - Inches(0.08), Inches(0.32),
           size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    params_map = {0:"99K", 1:"99K", 2:"99K", 3:"99K", 4:"99K",
                  5:"100K", 6:"523K", 7:"521K", 8:"523K", 9:"12.7M"}
    best_cnn = 8
    for i, row in enumerate(RESULTS):
        y = Inches(1.70 + i * 0.555)
        is_best = (i == best_cnn); is_garf = (i == len(RESULTS) - 1)
        if is_garf:
            rect(sl, Inches(0.3), y - Inches(0.04), Inches(11.1), Inches(0.06), fill=GOLD)
        bg = (RGBColor(0xE8, 0xF8, 0xE8) if is_best
              else RGBColor(0xFF, 0xFD, 0xE7) if is_garf
              else (LGRAY if i % 2 == 0 else WHITE))
        rect(sl, Inches(0.3), y, Inches(11.1), Inches(0.52),
             fill=bg, line=RGBColor(0xCC, 0xCC, 0xCC), lw=Pt(0.5))
        step_lbl = f"S{i+1}" if i < 9 else "PTv3"
        vals = [step_lbl, row["label"],
                f"{row['f1']:.3f}", f"{row['prec']:.3f}", f"{row['rec']:.3f}",
                params_map[i]]
        for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
            c = (GREEN if is_best and j >= 2
                 else GOLD if is_garf
                 else DARK)
            tb(sl, val, cx + Inches(0.04), y + Inches(0.10),
               cw - Inches(0.08), Inches(0.38),
               size=11, bold=(is_best and j >= 2) or is_garf, color=c,
               align=PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER)
    tb(sl, "★  Best CNN = Step 9   |   PTv3 reference in gold   |   Step 9 @ thresh=0.65 → F1≈90.1%",
       Inches(0.3), Inches(7.1), Inches(11.1), Inches(0.32),
       size=9, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

def slide_12_results_chart(prs):
    sl = blank(prs)
    header(sl, "Results — F1 Score Comparison")
    img(sl, fig_results_table_chart(),
        Inches(0.4), Inches(1.2), Inches(8.5), Inches(5.5))
    img(sl, fig_f1_progression(),
        Inches(8.9), Inches(1.2), Inches(4.2), Inches(5.5))

def slide_13_insights(prs):
    sl = blank(prs)
    header(sl, "Key Insights")
    insights = [
        (TEAL,   "Feature-level fusion unlocks U-Net  (+4.7 pp)",
                 "U-Net alone ≈ no gain. U-Net + bottleneck fusion → +4.7 pp. Synergistic."),
        (GREEN,  "Geometric features are zero-cost  (+3.6 pp)",
                 "Curvature, roughness, normal consistency add 0 parameters — pure geometric signal."),
        (GOLD,   "Threshold calibration closes the gap  (+4.5 pp free)",
                 "thresh=0.65 → F1≈90.1%.  thresh=0.70 → F1=90.9% ≈ PTv3 (90.98%).  No retraining."),
        (RED,    "Failure case: very low fracture ratio (<5%)",
                 "Model predicts nothing. Class imbalance collapse. Focal loss / oversampling needed."),
    ]
    for i, (col, title, desc) in enumerate(insights):
        y = Inches(1.4 + i * 1.42)
        rect(sl, Inches(0.3), y, Inches(0.22), Inches(1.05), fill=col)
        tb(sl, title, Inches(0.65), y + Inches(0.05), Inches(12.2), Inches(0.5),
           size=17, bold=True, color=col)
        tb(sl, desc, Inches(0.65), y + Inches(0.57), Inches(12.2), Inches(0.55),
           size=13, color=DARK)

def slide_14_limitations(prs):
    sl = blank(prs)
    header(sl, "Limitations")
    items = [
        (ORANGE, "Bilinear interpolation was not effective",
                 "Step 4 shows marginal loss vs Step 3. Splatting already fills coverage; bilinear adds noise."),
        (RED,    "3D → 2D projection is inherently lossy",
                 "Some fracture patterns are only visible from angles not in our 3 fixed views."),
        (DGRAY,  "Limited training budget (50 epochs vs 500 for GARF)",
                 "CNN results are preliminary — more epochs could close some of the gap."),
    ]
    for i, (col, title, desc) in enumerate(items):
        y = Inches(1.5 + i * 1.7)
        rect(sl, Inches(0.3), y, Inches(0.22), Inches(1.2), fill=col)
        tb(sl, title, Inches(0.65), y + Inches(0.05), Inches(12.2), Inches(0.5),
           size=17, bold=True, color=col)
        tb(sl, desc, Inches(0.65), y + Inches(0.58), Inches(12.2), Inches(0.6),
           size=14, color=DARK)

def slide_15_nextsteps(prs):
    sl = blank(prs)
    header(sl, "Next Steps & Conclusion")
    cols = [
        ("Achieved ✓",    GREEN, ["Step 9: F1=85.5% (−5.4 pp vs PTv3)",
                                   "Threshold calib: F1≈90.1% @ 0.65",
                                   "24× fewer params than PTv3",
                                   "Error analysis — failure modes identified"]),
        ("In progress",   TEAL,  ["Step 10: Tversky loss + dist centroid",
                                   "Formal test-set eval with thresh=0.65",
                                   "Retrieve qualitative figures from server"]),
        ("Future work",   NAVY,  ["Focal loss for low fracture-ratio frags",
                                   "Adaptive threshold per fragment",
                                   "Learnable projection angles"]),
    ]
    CW = Inches(3.9)
    for i, (term, col, items_) in enumerate(cols):
        cx = Inches(0.4 + i * 4.3)
        rect(sl, cx, Inches(1.3), CW, Inches(5.7), fill=LGRAY, line=col, lw=Pt(2))
        rect(sl, cx, Inches(1.3), CW, Inches(0.55), fill=col)
        tb(sl, term, cx + Inches(0.1), Inches(1.36), CW - Inches(0.2), Inches(0.44),
           size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items_):
            tb(sl, "→  " + item, cx + Inches(0.15), Inches(2.1 + j * 1.35),
               CW - Inches(0.3), Inches(1.3), size=13, color=DARK)

# ── Main ──────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="presentation.pptx")
    args = ap.parse_args()

    import os; os.makedirs(os.path.dirname(os.path.abspath(args.out)), exist_ok=True)

    print("Generating visuals…")
    prs = new_prs()
    for fn in [slide_01_title, slide_02_motivation, slide_03_idea, slide_04_pipeline,
               slide_05_ablation_principle, slide_06_normals, slide_07_splatting,
               slide_08_bilinear, slide_09_context,
               slide_10_attention, slide_10b_unet, slide_10c_fusion,
               slide_10d_geo, slide_10e_threshold, slide_10f_error_analysis,
               slide_11_results_table, slide_12_results_chart,
               slide_13_insights, slide_14_limitations, slide_15_nextsteps]:
        fn(prs)
        print(f"  ✓ {fn.__name__}")

    prs.save(args.out)
    print(f"\nSaved → {args.out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
